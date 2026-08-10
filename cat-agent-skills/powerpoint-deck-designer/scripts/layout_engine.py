"""Reusable slide layouts built directly with python-pptx shapes and charts."""

from __future__ import annotations

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from theme_engine import Theme
from text_fit import fit_font_size, truncate_bullets


BLANK_LAYOUT_INDEX = 6


class LayoutEngine:
    def __init__(self, prs, theme: Theme):
        self.prs = prs
        self.theme = theme
        self.slide_w = prs.slide_width
        self.slide_h = prs.slide_height

    def render(self, layout: str, spec: dict, slide_number: int, total: int) -> None:
        method = getattr(self, f"_layout_{layout.replace('-', '_')}", None)
        if method is None:
            method = self._layout_content
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[BLANK_LAYOUT_INDEX])
        self._paint_background(slide)
        method(slide, spec)
        if layout != "title":
            self._add_footer(slide, slide_number, total)
        self._add_speaker_notes(slide, spec.get("notes"))

    # -------- helpers --------
    def _paint_background(self, slide) -> None:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.slide_w, self.slide_h)
        bg.line.fill.background()
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.theme.rgb("background")
        bg.shadow.inherit = False

    def _add_footer(self, slide, num: int, total: int) -> None:
        tb = slide.shapes.add_textbox(
            Inches(0.4), self.slide_h - Inches(0.45),
            self.slide_w - Inches(0.8), Inches(0.35),
        )
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"{num} / {total}"
        run.font.name = self.theme.fonts["body"]
        run.font.size = Pt(10)
        run.font.color.rgb = self.theme.rgb("muted")

    def _add_speaker_notes(self, slide, notes) -> None:
        if not notes:
            return
        slide.notes_slide.notes_text_frame.text = notes

    def _styled_textbox(self, slide, left, top, width, height, text: str, *,
                        size_key: str, color_key: str = "text", bold: bool = False,
                        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                        font_key: str = "body") -> None:
        if not text:
            return
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        start = self.theme.sizes.get(size_key, 18)
        pt_size = fit_font_size(text, width, height, start_pt=start, min_pt=12)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = self.theme.fonts[font_key]
        run.font.size = Pt(pt_size)
        run.font.bold = bold
        run.font.color.rgb = self.theme.rgb(color_key)

    def _accent_bar(self, slide, left, top, width, height, color_key="accent") -> None:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme.rgb(color_key)

    # -------- layouts --------
    def _layout_title(self, slide, spec: dict) -> None:
        self._accent_bar(slide, 0, 0, Inches(0.35), self.slide_h, "primary")
        self._styled_textbox(slide, Inches(0.9), Inches(2.4),
                             self.slide_w - Inches(1.4), Inches(1.5),
                             spec.get("title", ""), size_key="title", bold=True,
                             color_key="primary", font_key="heading")
        self._styled_textbox(slide, Inches(0.9), Inches(4.0),
                             self.slide_w - Inches(1.4), Inches(1.0),
                             spec.get("subtitle", ""), size_key="subtitle",
                             color_key="text")
        footer_parts = [p for p in [spec.get("author"), spec.get("date")] if p]
        if footer_parts:
            self._styled_textbox(slide, Inches(0.9), self.slide_h - Inches(1.1),
                                 self.slide_w - Inches(1.4), Inches(0.5),
                                 " | ".join(footer_parts),
                                 size_key="caption", color_key="muted")

    def _layout_section(self, slide, spec: dict) -> None:
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5),
                                      self.slide_w, Inches(2.5))
        band.line.fill.background()
        band.fill.solid()
        band.fill.fore_color.rgb = self.theme.rgb("primary")
        self._styled_textbox(slide, Inches(0.6), Inches(2.9),
                             self.slide_w - Inches(1.2), Inches(1.7),
                             spec.get("title", ""), size_key="title",
                             bold=True, color_key="background",
                             align=PP_ALIGN.LEFT, font_key="heading")
        if spec.get("subtitle"):
            self._styled_textbox(slide, Inches(0.6), Inches(5.2),
                                 self.slide_w - Inches(1.2), Inches(0.6),
                                 spec["subtitle"], size_key="subtitle",
                                 color_key="muted")

    def _layout_content(self, slide, spec: dict) -> None:
        self._render_header(slide, spec)
        bullets = truncate_bullets(spec.get("bullets", []))
        if not bullets:
            self._styled_textbox(slide, Inches(0.6), Inches(1.9),
                                 self.slide_w - Inches(1.2), self.slide_h - Inches(2.6),
                                 spec.get("body", ""), size_key="body")
            return
        left = Inches(0.6)
        top = Inches(1.9)
        width = self.slide_w - Inches(1.2)
        height = self.slide_h - Inches(2.6)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        joined = "\n".join(bullets)
        pt_size = fit_font_size(joined, width, height,
                                start_pt=self.theme.sizes["body"], min_pt=14)
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = f"\u2022  {b}"
            run.font.name = self.theme.fonts["body"]
            run.font.size = Pt(pt_size)
            run.font.color.rgb = self.theme.rgb("text")

    def _layout_two_column(self, slide, spec: dict) -> None:
        self._render_header(slide, spec)
        col_w = int((self.slide_w - Inches(1.5)) / 2)
        top = Inches(1.9)
        height = self.slide_h - Inches(2.6)
        columns = spec.get("columns", [{"title": "", "body": ""},
                                       {"title": "", "body": ""}])
        for i, col in enumerate(columns[:2]):
            left = Inches(0.6) + i * (col_w + Inches(0.3))
            if col.get("title"):
                self._styled_textbox(slide, left, top, col_w, Inches(0.6),
                                     col["title"], size_key="heading", bold=True,
                                     color_key="primary", font_key="heading")
            body_top = top + Inches(0.7)
            body_h = height - Inches(0.7)
            bullets = truncate_bullets(col.get("bullets", []))
            text = "\n".join(f"\u2022  {b}" for b in bullets) if bullets else col.get("body", "")
            self._styled_textbox(slide, left, body_top, col_w, body_h,
                                 text, size_key="body")

    def _layout_cards(self, slide, spec: dict) -> None:
        self._render_header(slide, spec)
        cards = spec.get("cards", [])[:6]
        if not cards:
            return
        cols = 3 if len(cards) > 2 else len(cards)
        rows = -(-len(cards) // cols)
        gutter = Inches(0.25)
        area_left = Inches(0.6)
        area_top = Inches(1.9)
        area_w = self.slide_w - Inches(1.2)
        area_h = self.slide_h - Inches(2.6)
        card_w = int((area_w - gutter * (cols - 1)) / cols)
        card_h = int((area_h - gutter * (rows - 1)) / rows)
        for idx, card in enumerate(cards):
            r, c = divmod(idx, cols)
            left = area_left + c * (card_w + gutter)
            top = area_top + r * (card_h + gutter)
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         left, top, card_w, card_h)
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme.rgb("surface")
            box.line.color.rgb = self.theme.rgb("muted")
            box.line.width = Pt(0.75)
            box.shadow.inherit = False
            self._styled_textbox(slide, left + Inches(0.2), top + Inches(0.15),
                                 card_w - Inches(0.4), Inches(0.5),
                                 card.get("title", ""),
                                 size_key="heading", bold=True,
                                 color_key="primary", font_key="heading")
            self._styled_textbox(slide, left + Inches(0.2), top + Inches(0.75),
                                 card_w - Inches(0.4), card_h - Inches(0.9),
                                 card.get("body", ""), size_key="body")

    def _layout_agenda(self, slide, spec: dict) -> None:
        self._render_header(slide, {"title": spec.get("title", "Agenda"),
                                    "subtitle": spec.get("subtitle")})
        items = spec.get("items", [])[:8]
        top = Inches(1.9)
        row_h = min(Inches(0.7), int((self.slide_h - Inches(2.6)) / max(1, len(items))))
        for i, item in enumerate(items):
            y = top + i * row_h
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6),
                                            y + Inches(0.05), Inches(0.5), Inches(0.5))
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.theme.rgb("accent")
            circle.line.fill.background()
            circle.text_frame.text = str(i + 1)
            p = circle.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.bold = True
                r.font.color.rgb = self.theme.rgb("background")
                r.font.size = Pt(14)
                r.font.name = self.theme.fonts["heading"]
            self._styled_textbox(slide, Inches(1.3), y,
                                 self.slide_w - Inches(1.9), row_h,
                                 item if isinstance(item, str) else item.get("label", ""),
                                 size_key="body", bold=True, anchor=MSO_ANCHOR.MIDDLE)

    def _layout_quote(self, slide, spec: dict) -> None:
        self._styled_textbox(slide, Inches(1.5), Inches(1.5),
                             self.slide_w - Inches(3.0), Inches(0.8),
                             "\u201c", size_key="title", color_key="accent",
                             bold=True, font_key="heading")
        self._styled_textbox(slide, Inches(1.5), Inches(2.3),
                             self.slide_w - Inches(3.0), Inches(3.0),
                             spec.get("quote", ""), size_key="heading",
                             color_key="text", font_key="heading")
        if spec.get("attribution"):
            self._styled_textbox(slide, Inches(1.5), Inches(5.6),
                                 self.slide_w - Inches(3.0), Inches(0.6),
                                 f"\u2014 {spec['attribution']}",
                                 size_key="body", color_key="muted")

    def _layout_closing(self, slide, spec: dict) -> None:
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                      self.slide_w, self.slide_h)
        band.line.fill.background()
        band.fill.solid()
        band.fill.fore_color.rgb = self.theme.rgb("primary")
        self._styled_textbox(slide, Inches(0.8), Inches(2.4),
                             self.slide_w - Inches(1.6), Inches(1.4),
                             spec.get("title", "Thank you"), size_key="title",
                             bold=True, color_key="background", font_key="heading")
        self._styled_textbox(slide, Inches(0.8), Inches(4.0),
                             self.slide_w - Inches(1.6), Inches(1.0),
                             spec.get("subtitle", ""), size_key="subtitle",
                             color_key="background")
        if spec.get("contact"):
            self._styled_textbox(slide, Inches(0.8), self.slide_h - Inches(1.2),
                                 self.slide_w - Inches(1.6), Inches(0.6),
                                 spec["contact"], size_key="body",
                                 color_key="background")

    def _layout_chart(self, slide, spec: dict) -> None:
        self._render_header(slide, spec)
        chart_spec = spec.get("chart") or {}
        ctype = (chart_spec.get("type") or "column").lower()
        categories = list(chart_spec.get("categories") or [])
        series_specs = list(chart_spec.get("series") or [])
        show_legend = chart_spec.get("show_legend", True)
        show_labels = chart_spec.get("show_data_labels", True)
        stacked = chart_spec.get("stacked", False)

        if not categories or not series_specs:
            self._styled_textbox(slide, Inches(0.6), Inches(1.9),
                                 self.slide_w - Inches(1.2), self.slide_h - Inches(2.6),
                                 "(No chart data provided)",
                                 size_key="body", color_key="muted")
            return

        if ctype in ("pie", "donut") and len(series_specs) > 1:
            series_specs = series_specs[:1]

        stacked_norm = stacked
        if isinstance(stacked, str):
            stacked_norm = stacked.lower()
        is_100 = stacked_norm in ("100", "percent", "100%")
        is_stacked = bool(stacked) and not is_100

        chart_type_map = {
            ("column", False, False): XL_CHART_TYPE.COLUMN_CLUSTERED,
            ("column", True,  False): XL_CHART_TYPE.COLUMN_STACKED,
            ("column", False, True):  XL_CHART_TYPE.COLUMN_STACKED_100,
            ("bar",    False, False): XL_CHART_TYPE.BAR_CLUSTERED,
            ("bar",    True,  False): XL_CHART_TYPE.BAR_STACKED,
            ("bar",    False, True):  XL_CHART_TYPE.BAR_STACKED_100,
            ("line",   False, False): XL_CHART_TYPE.LINE_MARKERS,
            ("line",   True,  False): XL_CHART_TYPE.LINE_MARKERS,
            ("line",   False, True):  XL_CHART_TYPE.LINE_MARKERS,
            ("pie",    False, False): XL_CHART_TYPE.PIE,
            ("donut",  False, False): XL_CHART_TYPE.DOUGHNUT,
        }
        xl_type = chart_type_map.get((ctype, is_stacked, is_100))
        if xl_type is None:
            xl_type = chart_type_map.get((ctype, False, False),
                                         XL_CHART_TYPE.COLUMN_CLUSTERED)

        data = CategoryChartData()
        data.categories = [str(c) for c in categories]
        for s in series_specs:
            name = s.get("name", "Series")
            values = [float(v) if v is not None else 0.0
                      for v in (s.get("values") or [])]
            if len(values) < len(categories):
                values = values + [0.0] * (len(categories) - len(values))
            data.add_series(name, values[:len(categories)])

        left = Inches(0.6)
        top = Inches(1.9)
        width = self.slide_w - Inches(1.2)
        height = self.slide_h - Inches(2.6)

        gf = slide.shapes.add_chart(xl_type, left, top, width, height, data)
        chart = gf.chart
        chart.has_title = False

        chart.has_legend = bool(show_legend)
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.name = self.theme.fonts["body"]
            chart.legend.font.size = Pt(11)
            chart.legend.font.color.rgb = self.theme.rgb("text")

        palette = [self.theme.rgb(k) for k in ("primary", "secondary", "accent", "muted")]
        plot = chart.plots[0]

        effective_labels = show_labels
        if ctype == "line" and len(categories) > 6:
            effective_labels = False
        plot.has_data_labels = bool(effective_labels)
        if plot.has_data_labels:
            dl = plot.data_labels
            dl.font.size = Pt(10)
            dl.font.name = self.theme.fonts["body"]
            dl.font.color.rgb = self.theme.rgb("text")
            if ctype in ("pie", "donut"):
                dl.show_percentage = True
                dl.show_value = False
                dl.number_format = "0%"
            elif is_100:
                dl.show_value = True
                dl.number_format = "0%"
            else:
                dl.show_value = True
                try:
                    dl.position = XL_LABEL_POSITION.OUTSIDE_END
                except (ValueError, TypeError):
                    pass

        if ctype in ("pie", "donut"):
            series0 = plot.series[0]
            for i, _cat in enumerate(categories):
                point = series0.points[i]
                fill = point.format.fill
                fill.solid()
                fill.fore_color.rgb = palette[i % len(palette)]
                line = point.format.line
                line.color.rgb = self.theme.rgb("background")
                line.width = Pt(1.25)
        elif ctype == "line":
            for i, series in enumerate(plot.series):
                colour = palette[i % len(palette)]
                line = series.format.line
                line.color.rgb = colour
                line.width = Pt(2.25)
                try:
                    marker = series.marker
                    marker.format.fill.solid()
                    marker.format.fill.fore_color.rgb = colour
                    marker.format.line.color.rgb = colour
                except Exception:
                    pass
        else:
            for i, series in enumerate(plot.series):
                fill = series.format.fill
                fill.solid()
                fill.fore_color.rgb = palette[i % len(palette)]
                series.format.line.fill.background()

        if ctype in ("bar", "column", "line"):
            for axis in (chart.category_axis, chart.value_axis):
                axis.tick_labels.font.name = self.theme.fonts["body"]
                axis.tick_labels.font.size = Pt(10)
                axis.tick_labels.font.color.rgb = self.theme.rgb("text")

    # -------- shared header --------
    def _render_header(self, slide, spec: dict) -> None:
        self._accent_bar(slide, Inches(0.6), Inches(0.55),
                         Inches(0.4), Inches(0.08), "accent")
        self._styled_textbox(slide, Inches(0.6), Inches(0.7),
                             self.slide_w - Inches(1.2), Inches(0.8),
                             spec.get("title", ""), size_key="heading",
                             bold=True, color_key="primary", font_key="heading")
        if spec.get("subtitle"):
            self._styled_textbox(slide, Inches(0.6), Inches(1.45),
                                 self.slide_w - Inches(1.2), Inches(0.4),
                                 spec["subtitle"], size_key="caption",
                                 color_key="muted")
